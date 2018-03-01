import random
from pptx import Presentation
#import GUI as user_interface
study_sentences=[]
study_sentences_filled=[]
current_label=""
index=0
study_answer=[]
missed=[]
missed_filled=[]
valuesChecked=[]
study_sentences_filled_revised=[]
study_sentences_revised=[]
def convert_presentation(thisFile, thisBlankIntensity):
    global study_sentences_filled_revised
    global study_sentences_revised
    text_runs=[]#Instantiate the array that will store each sentence
    prs = Presentation(thisFile)#Load the presentation
    checkList=["the", "is", "a", "as", "of", "and", "or"]#Values in this list will not be converted into blanks
    file_split_string=""#will be used to join the words in the sentence back together
    for slide in prs.slides:#Go through each slide
        for shape in slide.shapes:#Go through each section
            if not shape.has_text_frame:#If it's not a text-frame
                continue#don't go through the rest of the program

            for paragraph in shape.text_frame.paragraphs:#Go through each part of the paragraph
                #Instance variables
                file_split=[]
                paragraph_text=""
                study_sentences_filled_text=""
                study_sentences_text=""
                for run in paragraph.runs:#Go through each sentence in the paragraph

                    text_runs.append(run.text)#Add that sentence to the array
                    study_sentences_filled_text+=(run.text.encode('ascii', 'ignore').decode('ascii'))#Add the correct answer in case of a quiz later         
                    file_split=run.text.split()#Split the sentence
                    for a in range(len(file_split)-1):#Go through each word in the sentence
                        for b in checkList:#Go through the checklist
                            if file_split[a]==b:#If the word is in the checklist
                                continue  #Definitely don't add a blank
                        #if (random.randint(0, round(len(file_split[a])/int(thisBlankIntensity)))==0):
                        if (random.randint(0, round((10/len(file_split[a])))*int(thisBlankIntensity))==0):#If at random this word is chosen
                        
                            file_split[a]="____"#Make the word a blank
                    file_split_string=" ".join(file_split)#Join the sentence back together
                    study_sentences_text+=file_split_string#Add the sentence to the one of the possible questions
                    paragraph_text+=file_split_string#Add the sentence to the paragraph
                study_sentences.append(study_sentences_text)#Add the sentence to the possible questions
                study_sentences_filled.append(study_sentences_filled_text)#Add a filled version of the sentence to the answer key
                paragraph.text=paragraph_text#Change the text to include the blanks
    print(study_sentences)
    study_sentences_revised=study_sentences
    study_sentences_filled_revised=study_sentences_filled
    newFileString=""
    
    for i in range(len(thisFile)-1):#Go through each character of the file name
        if (len(thisFile)-1)-i>4:#Make sure you don't include the powerpoint extension
            newFileString+=thisFile[i]#Add the file name without the extension to the string
    newFileString+="_blanks.pptx"#Add the extension
    print(newFileString)
    prs.save(str(newFileString))#Save the file
'''def question_check(i):
    global index
    if i in valuesChecked:

        print("Index selected: " + str(i) + " valuesChecked: " + str(valuesChecked))
        print("This value has already occured")
        index=random.randint(0, len(study_sentences)-1)
        question_check(index)
        '''
def study_blanks():
    #Instance variables
    global current_label
    global study_answer
    global study_sentences
    global index
    global study_sentences_revised
    study_answer=[]
    current_label=""
    current_label_local=""
    #End of instance variables
    if(len(study_sentences_revised)!=0):
        index=random.randint(0, len(study_sentences_revised)-1)
        print(index)
        '''print("******")
        print(index)
        print(valuesChecked)
        print("******")
        question_check(index)

        print(index)'''
        study_answer_split=study_sentences_filled_revised[index].split()

        #for i in study_sentences:
        index_split=study_sentences_revised[index].split()
        for a in range(0, len(index_split)-1):
            current_label_local+=index_split[a] + " "

            if index_split[a] =="____" and len(study_answer)<=3:
                study_answer.append(str(study_answer_split[a]))
        if(len(current_label)==0):
            current_label=study_sentences_revised[index]
            print(current_label_local)
            print(current_label)
            print(getCurrentLabel())
        print(study_answer)
        print(current_label)
def getCurrentLabel():
    print(current_label)
    return current_label

def getAnswer():
    return study_answer
